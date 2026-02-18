import os
import json
import pandas as pd
import pdfplumber
from docx import Document
from pptx import Presentation
from ebooklib import epub
import nltk

# Download NLTK resources (first run only)
nltk.download('punkt')
nltk.download('stopwords')
nltk.download('wordnet')


def extract_text_from_pdf(path):
    text = ""
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            if page.extract_text():
                text += page.extract_text() + "\n"
    return text.strip()

def extract_text_from_docx(path):
    doc = Document(path)
    return " ".join(p.text for p in doc.paragraphs).strip()

def extract_tables_from_excel(path):
    sheets = pd.read_excel(path, sheet_name=None)
    return sheets

def extract_text_from_ppt(path):
    prs = Presentation(path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text.strip()

def extract_text_from_epub(path):
    book = epub.read_epub(path)
    text = ""
    for item in book.get_items():
        if item.get_type() == epub.ITEM_DOCUMENT:
            text += item.get_content().decode("utf-8")
    return text.strip()


def normalize_to_json(file_type, content, source):
    return {
        "file_type": file_type,
        "content": content,
        "metadata": {
            "author": "Unknown",
            "date": "2025-02-18",
            "source": source
        }
    }



def preprocess_folder(input_folder, output_folder):
    os.makedirs(output_folder, exist_ok=True)

    for file in os.listdir(input_folder):
        file_path = os.path.join(input_folder, file)

        try:
            if file.endswith(".pdf"):
                content = extract_text_from_pdf(file_path)
                data = normalize_to_json("PDF", content, file)

            elif file.endswith(".docx"):
                content = extract_text_from_docx(file_path)
                data = normalize_to_json("Word", content, file)

            elif file.endswith(".xlsx"):
                tables = extract_tables_from_excel(file_path)
                excel_content = {}

                for sheet_name, df in tables.items():
                    df = df.astype(str)  
                    excel_content[sheet_name] = df.to_dict(orient="records")
                
                data = normalize_to_json("Excel", excel_content, file)


            elif file.endswith(".pptx"):
                content = extract_text_from_ppt(file_path)
                data = normalize_to_json("PowerPoint", content, file)

            elif file.endswith(".epub"):
                content = extract_text_from_epub(file_path)
                data = normalize_to_json("EPUB", content, file)

            else:
                continue

            output_file = os.path.join(
                output_folder,
                file.replace(".", "_") + ".json"
            )

            with open(output_file, "w", encoding="utf-8") as f:
                json.dump(data, f, indent=4, ensure_ascii=False)

            print(f"[✔] Processed: {file}")

        except Exception as e:
            print(f"[✘] Error processing {file}: {e}")



if __name__ == "__main__":
    preprocess_folder("input_documents", "output")
